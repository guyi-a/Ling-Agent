"""
文件工具 - 提供文件系统的读写操作

工作区路径规则：
  - 所有操作限制在 WORKSPACE_ROOT/{project_id}_{slug}/ 内
  - 相对路径自动解析为工作区内的路径
  - 绝对路径必须在工作区内，否则拒绝访问

目录结构：
  WORKSPACE_ROOT/
  └── {project_id}_{slug}/
      ├── uploads/    # 用户上传的原始文件
      └── ...         # Agent 生成的文件直接放根目录
"""
import logging
from pathlib import Path
from typing import Type, Optional

from langchain.tools import BaseTool
from pydantic import BaseModel, Field

from app.agent.tools._ctx import get_session_id
from app.core.config import settings

logger = logging.getLogger(__name__)


def get_session_workspace(session_id: str, *, ensure: bool = True) -> Path:
    """获取当前 session 的工作区目录。

    优先使用项目 slug 目录（物化项目），否则回退到 session_id 目录。
    ensure=True 时自动创建目录（工具写入时使用）。
    """
    from app.database.session import sync_session_factory
    from app.models.session import Session
    from app.models.project import Project
    from sqlalchemy import select

    try:
        with sync_session_factory() as db:
            result = db.execute(
                select(Session).where(Session.session_id == session_id)
            )
            session = result.scalars().first()
            if session and session.project_id:
                proj_result = db.execute(
                    select(Project).where(Project.id == session.project_id)
                )
                project = proj_result.scalars().first()
                if project and project.slug:
                    workspace = Path(settings.WORKSPACE_ROOT).resolve() / f"{project.id}_{project.slug}"
                    if ensure:
                        workspace.mkdir(parents=True, exist_ok=True)
                        (workspace / "uploads").mkdir(exist_ok=True)
                    return workspace
    except Exception:
        pass

    workspace = Path(settings.WORKSPACE_ROOT).resolve() / session_id
    if ensure:
        workspace.mkdir(parents=True, exist_ok=True)
        (workspace / "uploads").mkdir(exist_ok=True)
    return workspace


def resolve_path(path: str, session_id: Optional[str] = None) -> Path:
    """
    解析文件路径：
    - 相对路径 → WORKSPACE_ROOT/{session_id}/{path}
    - 绝对路径 → 原路径（必须在 WORKSPACE_ROOT 内）
    """
    p = Path(path).expanduser()

    if not p.is_absolute():
        if session_id:
            base = get_session_workspace(session_id)
        else:
            base = Path(settings.WORKSPACE_ROOT)
            base.mkdir(parents=True, exist_ok=True)
        p = (base / path).resolve()
    else:
        p = p.resolve()

    # 安全检查：路径必须在 WORKSPACE_ROOT 内
    workspace_root = Path(settings.WORKSPACE_ROOT).resolve()
    try:
        p.relative_to(workspace_root)
    except ValueError:
        raise PermissionError(
            f"Access denied: path '{p}' is outside workspace '{workspace_root}'"
        )

    return p


class _ReadFileInput(BaseModel):
    path: str = Field(description="文件路径（相对于工作区的相对路径，如 uploads/data.csv 或 report.md）")


class _WriteFileInput(BaseModel):
    path: str = Field(description="文件路径（相对于工作区的相对路径，如 report.md、chart.png）")
    content: str = Field(description="要写入的内容")


class _ListDirInput(BaseModel):
    path: str = Field(default=".", description="目录路径（默认为 session 工作区根目录）")


class ReadFileTool(BaseTool):
    """读取工作区内的文件内容（支持纯文本和文档格式）"""
    name: str = "read_file"
    description: str = (
        "Read file contents from the workspace. "
        "Supports:\n"
        "- Text files: .txt, .md, .csv, .py, .js, .json, etc.\n"
        "- Documents: .docx (Word), .pdf (PDF), .pptx (PowerPoint)\n"
        "For documents, automatically extracts text, tables, and images.\n"
        "Use relative paths like 'uploads/report.pdf' or 'data.csv'."
    )
    args_schema: Type[BaseModel] = _ReadFileInput
    current_session_id: Optional[str] = None

    def _run(self, path: str) -> str:
        try:
            p = resolve_path(path, get_session_id())
            if not p.exists():
                return f"Error: File not found: {path}"
            if not p.is_file():
                return f"Error: Path is not a file: {path}"

            # 根据扩展名选择处理方式
            ext = p.suffix.lower()

            if ext == '.docx':
                return self._read_docx(p)
            elif ext == '.pdf':
                return self._read_pdf(p)
            elif ext == '.pptx':
                return self._read_pptx(p)
            else:
                # 纯文本读取
                content = p.read_text(encoding="utf-8")
                logger.info(f"📖 Read text file: {p} ({len(content)} chars)")
                return content

        except PermissionError as e:
            return f"Error: {e}"
        except UnicodeDecodeError:
            return f"Error: Cannot read binary file '{path}' as text. Use a document tool or convert it first."
        except Exception as e:
            logger.error(f"Error reading file: {e}")
            return f"Error reading file '{path}': {e}"

    def _read_docx(self, path: Path) -> str:
        """读取 Word 文档"""
        try:
            from docx import Document

            doc = Document(str(path))

            # 提取文本
            text = "\n\n".join(p.text.strip() for p in doc.paragraphs if p.text.strip())

            # 提取表格
            tables = []
            for table in doc.tables:
                rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
                if rows:
                    tables.append(rows)

            logger.info(f"📄 Read Word: {len(text)} chars, {len(tables)} tables")

            # 格式化返回结果
            output = f"📄 Document: {path.name}\n"
            output += f"📊 Pages: {len(doc.paragraphs)}\n"
            output += f"📝 Text length: {len(text)} characters\n"

            if tables:
                output += f"📋 Tables: {len(tables)}\n"

            output += f"\n--- Content ---\n{text}"
            return output

        except ImportError:
            return "Error: python-docx library not installed. Run: pip install python-docx"
        except Exception as e:
            return f"Error reading Word document: {e}"

    def _read_pdf(self, path: Path) -> str:
        """读取 PDF 文档（支持 OCR fallback）"""
        try:
            import fitz  # PyMuPDF

            doc = fitz.open(str(path))

            # 提取文本
            text_parts = []
            ocr_pages = []
            for page_num, page in enumerate(doc, 1):
                page_text = page.get_text().strip()

                # 如果文本提取成功且内容足够（>20 字符），直接使用
                if page_text and len(page_text) > 20:
                    text_parts.append(f"--- Page {page_num} ---\n{page_text}")
                else:
                    # 尝试 OCR
                    ocr_text = self._try_ocr_page(page)
                    if ocr_text:
                        text_parts.append(f"--- Page {page_num} (OCR) ---\n{ocr_text}")
                        ocr_pages.append(page_num)
                    elif page_text:
                        # 有少量文本但不足 20 字符，仍然保留
                        text_parts.append(f"--- Page {page_num} ---\n{page_text}")

            text = "\n\n".join(text_parts)
            page_count = doc.page_count

            logger.info(f"📕 Read PDF: {page_count} pages, {len(text)} chars, OCR used on {len(ocr_pages)} pages")

            doc.close()

            output = f"📄 Document: {path.name}\n"
            output += f"📊 Pages: {page_count}\n"
            output += f"📝 Text length: {len(text)} characters\n"
            if ocr_pages:
                output += f"🔍 OCR used on pages: {', '.join(map(str, ocr_pages))}\n"
            output += f"\n--- Content ---\n{text}"
            return output

        except ImportError:
            return "Error: PyMuPDF library not installed. Run: pip install PyMuPDF"
        except Exception as e:
            return f"Error reading PDF document: {e}"

    def _try_ocr_page(self, page) -> str:
        """对单页 PDF 执行 OCR（使用 RapidOCR）"""
        try:
            import cv2
            import numpy as np
            from rapidocr_onnxruntime import RapidOCR
            import fitz

            # 将页面转为图像（2倍分辨率提高 OCR 准确度）
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)

            # 转换为 numpy array
            channels = 4 if pix.alpha else 3
            image = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.height, pix.width, channels)
            if channels == 4:
                image = cv2.cvtColor(image, cv2.COLOR_BGRA2BGR)

            # 执行 OCR
            engine = RapidOCR()
            result, _ = engine(image)

            # 提取文本行
            lines = []
            for item in result or []:
                if len(item) >= 2:
                    text = str(item[1]).strip()
                    if text:
                        lines.append(text)

            return '\n'.join(lines) if lines else ""

        except ImportError:
            logger.warning("OCR dependencies not available (need opencv-python, rapidocr-onnxruntime)")
            return ""
        except Exception as e:
            logger.warning(f"OCR failed: {e}")
            return ""

    def _read_pptx(self, path: Path) -> str:
        """读取 PowerPoint 文档"""
        try:
            from pptx import Presentation

            prs = Presentation(str(path))

            # 提取每页文本
            text_parts = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text.strip())
                if slide_texts:
                    text_parts.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_texts))

            text = "\n\n".join(text_parts)

            logger.info(f"📊 Read PPTX: {len(prs.slides)} slides, {len(text)} chars")

            output = f"📄 Document: {path.name}\n"
            output += f"📊 Slides: {len(prs.slides)}\n"
            output += f"📝 Text length: {len(text)} characters\n"
            output += f"\n--- Content ---\n{text}"
            return output

        except ImportError:
            return "Error: python-pptx library not installed. Run: pip install python-pptx"
        except Exception as e:
            return f"Error reading PowerPoint document: {e}"

    async def _arun(self, path: str) -> str:
        return self._run(path)


class WriteFileTool(BaseTool):
    """在工作区内写入文件（不存在则创建，存在则覆盖）"""
    name: str = "write_file"
    description: str = (
        "Write content to a file in the workspace. "
        "Use relative paths like 'report.md'. "
        "Creates parent directories automatically."
    )
    args_schema: Type[BaseModel] = _WriteFileInput
    current_session_id: Optional[str] = None

    def _run(self, path: str, content: str) -> str:
        try:
            p = resolve_path(path, get_session_id())
            p.parent.mkdir(parents=True, exist_ok=True)
            p.write_text(content, encoding="utf-8")
            logger.info(f"✏️  Wrote file: {p} ({len(content)} chars)")
            return f"Successfully wrote {len(content)} characters to {p}"
        except PermissionError as e:
            return f"Error: {e}"
        except Exception as e:
            return f"Error writing file '{path}': {e}"

    async def _arun(self, path: str, content: str) -> str:
        return self._run(path, content)


class _EditFileInput(BaseModel):
    path: str = Field(description="文件路径（相对于工作区的相对路径，如 index.html）")
    old_string: str = Field(description="要被替换的原始文本（必须与文件中的内容精确匹配）")
    new_string: str = Field(description="替换后的新文本")


class EditFileTool(BaseTool):
    """在工作区内对文件做局部替换（精确匹配 old_string → 替换为 new_string）"""
    name: str = "edit_file"
    description: str = (
        "Edit a file by replacing a specific string with new content. "
        "Much more efficient than rewriting the entire file — use this for small changes. "
        "The old_string must match EXACTLY (including whitespace and indentation). "
        "If old_string appears multiple times, only the FIRST occurrence is replaced. "
        "Use relative paths like 'index.html'."
    )
    args_schema: Type[BaseModel] = _EditFileInput
    current_session_id: Optional[str] = None

    def _run(self, path: str, old_string: str, new_string: str) -> str:
        try:
            p = resolve_path(path, get_session_id())
            if not p.exists():
                return f"Error: File not found: {path}"
            if not p.is_file():
                return f"Error: Path is not a file: {path}"

            content = p.read_text(encoding="utf-8")

            if old_string not in content:
                return (
                    f"Error: old_string not found in {path}. "
                    "Make sure it matches exactly (including whitespace and indentation)."
                )

            if old_string == new_string:
                return "Error: old_string and new_string are identical, nothing to change."

            new_content = content.replace(old_string, new_string, 1)
            p.write_text(new_content, encoding="utf-8")

            logger.info(
                f"✏️  Edited file: {p} "
                f"(replaced {len(old_string)} chars → {len(new_string)} chars)"
            )
            return (
                f"Successfully edited {path}: "
                f"replaced {len(old_string)} chars with {len(new_string)} chars."
            )
        except PermissionError as e:
            return f"Error: {e}"
        except UnicodeDecodeError:
            return f"Error: Cannot edit binary file '{path}'."
        except Exception as e:
            return f"Error editing file '{path}': {e}"

    async def _arun(self, path: str, old_string: str, new_string: str) -> str:
        return self._run(path, old_string, new_string)


class ListDirTool(BaseTool):
    """列出工作区目录内容"""
    name: str = "list_dir"
    description: str = (
        "List files and directories in the workspace. "
        "Default lists the session root. Use 'uploads' for user-uploaded files."
    )
    args_schema: Type[BaseModel] = _ListDirInput
    current_session_id: Optional[str] = None

    def _run(self, path: str = ".") -> str:
        try:
            p = resolve_path(path, get_session_id())
            if not p.exists():
                return f"Directory is empty or not found: {path}"
            if not p.is_dir():
                return f"Error: Path is not a directory: {path}"

            entries = sorted(p.iterdir(), key=lambda x: (x.is_file(), x.name))
            if not entries:
                return f"Directory is empty: {p}"

            lines = []
            for entry in entries:
                prefix = "📄 " if entry.is_file() else "📁 "
                size = f" ({entry.stat().st_size} bytes)" if entry.is_file() else ""
                lines.append(f"{prefix}{entry.name}{size}")

            logger.info(f"📂 Listed dir: {p} ({len(entries)} entries)")
            return f"Contents of {p}:\n" + "\n".join(lines)
        except PermissionError as e:
            return f"Error: {e}"
        except Exception as e:
            return f"Error listing directory '{path}': {e}"

    async def _arun(self, path: str = ".") -> str:
        return self._run(path)
